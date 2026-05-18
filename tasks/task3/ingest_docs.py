import re
import psycopg2
from pathlib import Path
from docx import Document
from pptx import Presentation
from pgvector.psycopg2 import register_vector
from sentence_transformers import SentenceTransformer

from db import DB_URL

DOCS_DIR = Path(__file__).resolve().parents[2] / "data" / "docs"
CHUNK_SIZE    = 400   # words per chunk
CHUNK_OVERLAP = 50    # overlap so we don't cut sentences in half

embedder = SentenceTransformer("all-MiniLM-L6-v2")


def extract_text_docx(path: Path) -> list[tuple[str, str]]:
    """Returns list of (section_heading, paragraph_text) tuples."""
    doc = Document(path)
    sections = []
    current_section = "General"
    buffer = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            if buffer:
                sections.append((current_section, " ".join(buffer)))
                buffer = []
            current_section = text
        else:
            buffer.append(text)

    if buffer:
        sections.append((current_section, " ".join(buffer)))

    # also pull text out of tables
    for table in doc.tables:
        table_text = []
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    table_text.append(cell.text.strip())
        if table_text:
            sections.append(("Table", " ".join(table_text)))

    return sections


def extract_text_pptx(path: Path) -> list[tuple[str, str]]:
    prs = Presentation(path)
    sections = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            sections.append(("Slide", " ".join(texts)))
    return sections


def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i + size])
        if chunk.strip():
            chunks.append(chunk)
        i += size - overlap
    return chunks


def parse_filename(path: Path) -> tuple[str, str, str]:
    name   = path.stem
    region = path.parent.name

    # date like "April 2018"
    date_match = re.search(r'(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}', name)
    doc_date = date_match.group(0) if date_match else ""

    # station name sits between "Plan - " and " Issue" in the filename
    station_match = re.search(r'Plan\s*-\s*(.+?)\s+Issue', name)
    if station_match:
        station = station_match.group(1).strip()
    else:
        station = name.replace("Station Disruption Plan", "").strip(" -_")

    return station, region, doc_date


def ingest_file(path: Path, cursor) -> int:
    cursor.execute("SELECT 1 FROM documents WHERE source_file = %s LIMIT 1", (path.name,))
    if cursor.fetchone():
        print(f"  already ingested, skipping")
        return 0

    station, region, doc_date = parse_filename(path)
    suffix = path.suffix.lower()

    if suffix in (".docx", ".doc", ".docm"):
        try:
            sections = extract_text_docx(path)
        except Exception as e:
            print(f"  Skipped {path.name}: {e}")
            return 0
    elif suffix in (".pptx", ".pptm"):
        try:
            sections = extract_text_pptx(path)
        except Exception as e:
            print(f"  Skipped {path.name}: {e}")
            return 0
    else:
        return 0

    count = 0
    for section, text in sections:
        for chunk in chunk_text(text):
            embedding = embedder.encode(chunk).tolist()
            cursor.execute(
                """
                INSERT INTO documents (station, region, doc_date, source_file, section, chunk_text, embedding)
                VALUES (%s, %s, %s, %s, %s, %s, %s::vector)
                """,
                (station, region, doc_date, path.name, section, chunk, embedding)
            )
            count += 1

    return count


def main():
    conn   = psycopg2.connect(DB_URL)
    register_vector(conn)
    cursor = conn.cursor()

    files = list(DOCS_DIR.rglob("*"))
    files = [f for f in files if f.suffix.lower() in (".docx", ".doc", ".docm", ".pptx", ".pptm")]

    print(f"Found {len(files)} files")
    total = 0

    try:
        for i, path in enumerate(files, 1):
            print(f"[{i}/{len(files)}] {path.name}")
            count = ingest_file(path, cursor)
            total += count
            print(f"  → {count} chunks")
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        cursor.close()
        conn.close()

    print(f"\nDone. {total} chunks stored.")


if __name__ == "__main__":
    main()
